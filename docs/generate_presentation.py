"""NullRefer Music — 발표자료 생성 (이미지 포함)"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ASSETS = Path(__file__).resolve().parent / "presentation_assets"
OUT = Path(__file__).resolve().parent / "NullReferMusic_Presentation.pptx"

PRIMARY = RGBColor(0, 102, 204)
INK = RGBColor(30, 30, 30)
MUTED = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    img = ASSETS / "slide1-app-overview.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(5.2), Inches(1.2), width=Inches(4.3))

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(4.5), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "NullRefer Music"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = "YouTube 음원 검색 · 차트 · 메타데이터 다운로드"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(220, 235, 255)
    p2.space_before = Pt(10)

    p3 = tf.add_paragraph()
    p3.text = "웹 + Android/iOS 릴리스 · AI(Cursor) 보조 개발"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(200, 220, 245)
    p3.space_before = Pt(8)


def add_slide_with_image(
    prs: Presentation,
    title: str,
    bullets: list[str],
    image_name: str,
    footer: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.4), Inches(5.8), Inches(0.65))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = INK

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(5.5), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15) if not line.startswith("  ") else Pt(13)
        p.font.color.rgb = MUTED if line.startswith("  ") else INK
        p.level = 1 if line.startswith("  ") else 0
        p.space_after = Pt(8)

    img_path = ASSETS / image_name
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(6.15), Inches(1.05), width=Inches(3.5))

    if footer:
        fb = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(9), Inches(0.35))
        fp = fb.text_frame.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(10)
        fp.font.color.rgb = MUTED


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_slide_with_image(
        prs,
        "프로그램 소개",
        [
            "NullRefer Music — YouTube 기반 음원 검색·다운로드 앱",
            "차트(Last.fm·Spotify 등)에서 곡을 고르거나 직접 검색",
            "다운로드 시 트랙 정보(메타데이터) 편집 후 파일로 저장",
            "동일 코드베이스로 웹 브라우저 + Android APK + iOS IPA 제공",
        ],
        "slide1-app-overview.png",
    )

    add_slide_with_image(
        prs,
        "개발 스택 · 환경",
        [
            "프론트: Expo(React Native) + TypeScript",
            "웹: 브라우저 + 로컬 Spring Boot API(8787)",
            "Android 릴리스: Chaquopy + yt-dlp·FFmpeg 온디바이스",
            "iOS: Innertube(youtubei.js) 스트림 + 기기 저장",
            "도구: yt-dlp, ffmpeg, Metro 번들러, Gradle",
        ],
        "slide2-dev-stack.png",
    )

    add_slide_with_image(
        prs,
        "배포 · Expo Go",
        [
            "최종 산출: 웹(브라우저) + Android APK + iOS IPA",
            "모바일 UI·기능 테스트: Expo Go + PC Metro(QR)",
            "  Expo Go = 폰에 설치하는 Expo 공식 샌드박스 앱",
            "  PC에서 JS 수정 → QR로 폰에 즉시 반영(빠른 테스트)",
            "릴리스 APK/IPA는 PC 서버 없이 단독 실행",
        ],
        "slide3-release-expo.png",
    )

    add_slide_with_image(
        prs,
        "주요 기능 · API · 다운로드",
        [
            "YouTube: Data API(검색) / Innertube(앱·폴백)",
            "Last.fm: 차트·검색·트랙 메타 API",
            "Spotify: 공식 API·Charts 세션(차트/기간별)",
            "다운로드 흐름:",
            "  1) yt-dlp로 오디오 추출 (웹=백엔드, Android=기기)",
            "  2) Last.fm 등으로 메타 보강 (병렬 prefetch)",
            "  3) ffmpeg로 MP3/m4a 변환 + ID3·커버 임베드",
        ],
        "slide4-download-api.png",
    )

    add_slide_with_image(
        prs,
        "AI(Cursor) 활용",
        [
            "구현: 요구사항·목업 → 코드 탐색·수정안 → 최소 diff 적용",
            "  UI·다운로드 파이프라인·토큰 갱신 등 반복 작업 가속",
            "테스트: tsc 타입 검사, 웹·Expo Go·APK 시나리오 점검",
            "  오류 로그 기반 원인 추적·예외 처리 보완",
            "사람: UX 의도·API 키·릴리스 빌드·최종 동작 확인",
        ],
        "slide5-ai-dev.png",
        footer="NullRefer Music — 2026",
    )

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
